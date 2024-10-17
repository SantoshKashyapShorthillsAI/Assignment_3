from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()

# List of hyperlinks to add to different slides
hyperlinks = [
    ("PowerPoint with Links, Text, Table, and Images", "This is a Python-generated PowerPoint with multiple elements.", "https://www.python.org", "Python's official website"),
    ("Table Example", "Here is a simple table example:", None, None),  # No hyperlink for this slide
    ("Image Example", "Check out this image:", None, None),  # No hyperlink for this slide
]

# Add slides for hyperlinks
for title_text, body_text, link, link_text in hyperlinks:
    slide_layout = presentation.slide_layouts[5]  # 5 is a blank slide
    slide = presentation.slides.add_slide(slide_layout)

    # Add Title Text
    title = slide.shapes.title
    title.text = title_text

    # Add body text box
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = body_text

    # Add hyperlink if provided
    if link and link_text:
        p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = f"Click here to visit {link_text}"
        run.hyperlink.address = link

    # Animation placeholder: Title and text will be animated when added later
    # This can be done in PowerPoint manually by going to "Animations" -> Choose Animation.

# Add a slide for a table
slide_layout = presentation.slide_layouts[5]
slide2 = presentation.slides.add_slide(slide_layout)

# Add a title for the table slide
title = slide2.shapes.title
title.text = "Table Example"

# Create a table
rows = 2
cols = 2
left = Inches(2)
top = Inches(2)
width = Inches(4)
height = Inches(1)
table = slide2.shapes.add_table(rows, cols, left, top, width, height).table

# Set column headers
table.cell(0, 0).text = 'Name'
table.cell(0, 1).text = 'Age'

# Set row content
table.cell(1, 0).text = 'John'
table.cell(1, 1).text = '30'

# Add a slide for an image
slide3 = presentation.slides.add_slide(slide_layout)

# Add a title for the image slide
title = slide3.shapes.title
title.text = "Image Example"

# Add an image (Make sure you have an image in the same directory as the script)
img_path = "/home/shtlp_0103/Assignment_3/Documents/apple.jpeg"  # Update with your image path
left = Inches(1)
top = Inches(2)
height = Inches(3)
slide3.shapes.add_picture(img_path, left, top, height=height)

# Animation placeholder: Image will be animated later in PowerPoint manually

# Save the presentation
presentation.save('animations.pptx')
print("Presentation created successfully! Please add animations manually in PowerPoint.")
