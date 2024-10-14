from abc import ABC, abstractmethod
import fitz
import docx
from pptx import Presentation
import os
import csv
import pdfplumber
import mysql.connector


class FileLoader(ABC):
    @abstractmethod
    def validate(self):
        """Check if the file is valid."""
        pass

    @abstractmethod
    def load(self):
        """Load the file content."""
        pass


class PDFLoader(FileLoader):
    def __init__(self, file_path):
        self.file_path = file_path
        self.doc = None

    def validate(self):
        return self.file_path.endswith('.pdf')

    def load(self):
        if not self.validate():
            raise ValueError("Invalid file format for PDF.")
        self.doc = fitz.open(self.file_path)
        return self.doc


class DOCXLoader(FileLoader):
    def __init__(self, file_path):
        self.file_path = file_path
        self.doc = None

    def validate(self):
        return self.file_path.endswith('.docx')

    def load(self):
        if not self.validate():
            raise ValueError("Invalid file format for DOCX.")
        self.doc = docx.Document(self.file_path)
        return self.doc


class PPTLoader(FileLoader):
    def __init__(self, file_path):
        self.file_path = file_path
        self.presentation = None

    def validate(self):
        return self.file_path.endswith('.pptx')

    def load(self):
        if not self.validate():
            raise ValueError("Invalid file format for PPTX.")
        self.presentation = Presentation(self.file_path)
        return self.presentation

class DataExtractor:
    def __init__(self, file_loader: FileLoader):
        self.file_loader = file_loader
        self.file_loader.load()

    def extract_text(self):
        """Extracts text with metadata like page number and font details."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_text()
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_text()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_text()
        else:
            raise ValueError("Unsupported file type for text extraction.")

    def _extract_pdf_text(self):
        text_data = []
        for page_num in range(len(self.file_loader.doc)):
            page = self.file_loader.doc.load_page(page_num)
            text = page.get_text("text")
            text_data.append({
                "page_number": page_num + 1,
                "text": text
            })
        return text_data

    def _extract_docx_text(self):
        text_data = []
        for para in self.file_loader.doc.paragraphs:
            text_data.append({
                "text": para.text,
                "style": para.style.name
            })
        return text_data

    def _extract_ppt_text(self):
        text_data = []
        for slide_num, slide in enumerate(self.file_loader.presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            # Join the slide text list into a single string separated by newlines
            text_data.append({
                "slide_number": slide_num + 1,
                "text": "\n".join(slide_text)  # Join text list into a single string
            })
        return text_data

    def extract_links(self):
        """Extracts hyperlinks with metadata."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_links()
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_links()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_links()
        else:
            raise ValueError("Unsupported file type for link extraction.")

    def _extract_pdf_links(self):
        link_data = []
        for page_num in range(len(self.file_loader.doc)):
            page = self.file_loader.doc.load_page(page_num)
            links = page.get_links()
            for link in links:
                link_data.append({
                    "page_number": page_num + 1,
                    "url": link.get('uri')
                })
        return link_data

    def _extract_docx_links(self):
        link_data = []
        for para in self.file_loader.doc.paragraphs:
            for run in para.runs:
                if run.font.color and run.text.startswith('http'):
                    link_data.append({
                        "url": run.text,
                        "style": para.style.name
                    })
        return link_data

    def _extract_ppt_links(self):
        link_data = []
        for slide_num, slide in enumerate(self.file_loader.presentation.slides):
            for shape in slide.shapes:
                # Check if the shape contains text and has a hyperlink attribute
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink and run.hyperlink.address:
                                link_data.append({
                                    "slide_number": slide_num + 1,
                                    "url": run.hyperlink.address
                                })
                # In case the shape has a hyperlink directly (without being in the text frame)
                elif hasattr(shape, "hyperlink") and shape.hyperlink.address:
                    link_data.append({
                        "slide_number": slide_num + 1,
                        "url": shape.hyperlink.address
                    })
        
        return link_data


    def extract_images(self):
        """Extract images from the document."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_images()
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_images()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_images()
        else:
            raise ValueError("Unsupported file type for image extraction.")

    def _extract_pdf_images(self):
        image_data = []
        for page_num in range(len(self.file_loader.doc)):
            page = self.file_loader.doc.load_page(page_num)
            images = page.get_images(full=True)
            for img_index, img in enumerate(images):
                xref = img[0]
                image = self.file_loader.doc.extract_image(xref)
                img_bytes = image["image"]
                img_extension = image["ext"]
                image_data.append({
                    "page_number": page_num + 1,
                    "image_data": img_bytes,
                    "image_extension": img_extension
                })
        return image_data

    def _extract_docx_images(self):
        image_data = []
        for rel in self.file_loader.doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data.append({
                    "image_data": rel.target_part.blob,
                    "image_extension": rel.target_part.content_type.split('/')[-1]
                })
        return image_data

    def _extract_ppt_images(self):
        image_data = []
        for slide_num, slide in enumerate(self.file_loader.presentation.slides):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    img_bytes = shape.image.blob
                    image_extension = shape.image.content_type.split('/')[-1]
                    image_data.append({
                        "slide_number": slide_num + 1,
                        "image_data": img_bytes,
                        "image_extension": image_extension
                    })
        return image_data

    def extract_tables(self):
        """Extract tables from the document."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_tables_with_plumber()  # Using 
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_tables()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_tables()  # Tables are rare in PPT, but handle if any.
        else:
            raise ValueError("Unsupported file type for table extraction.")

    def _extract_pdf_tables_with_plumber(self):
        """Extract tables from PDF using pdfplumber."""
        table_data = []
        with pdfplumber.open(self.file_loader.file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                tables = page.extract_tables()  # Extract tables from each page
                for table in tables:
                    table_data.append({
                        "page_number": page_num + 1,
                        "table": table
                    })
        return table_data
    
    def _extract_docx_tables(self):
        table_data = []
        for table in self.file_loader.doc.tables:
            table_rows = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                table_rows.append(row_data)
            table_data.append({
                "table": table_rows
            })
        return table_data

    def _extract_ppt_tables(self):
        table_data = []
        for slide_num, slide in enumerate(self.file_loader.presentation.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_rows.append(row_data)
                    table_data.append({
                        "slide_number": slide_num + 1,
                        "table": table_rows
                    })
        return table_data

class Storage(ABC):
    """Abstract class for storing extracted data."""

    @abstractmethod
    def save_text(self, text_data):
        """Save extracted text."""
        pass

    @abstractmethod
    def save_images(self, images_data):
        """Save extracted images."""
        pass

    @abstractmethod
    def save_tables(self, tables_data):
        """Save extracted tables."""
        pass

    @abstractmethod
    def save_links(self, links_data):
        """Save extracted links."""
        pass


class FileStorage(Storage):
    """Concrete class for storing extracted data to files."""

    def __init__(self, output_directory):
        """Initialize FileStorage with an output directory."""
        self.output_directory = output_directory
        if not os.path.exists(self.output_directory):
            os.makedirs(self.output_directory)

    def save_text(self, text_data):
        """Save extracted text to a text file."""
        with open(os.path.join(self.output_directory, 'extracted_text.txt'), 'w') as f:
            for entry in text_data:
                f.write(f"{entry}\n")

    def save_images(self, images_data):
        """Save extracted images to the output directory."""
        for i, image in enumerate(images_data):
            image_extension = image.get("image_extension", "png")  # Default to PNG if not provided
            image_path = os.path.join(self.output_directory, f'image_{i}.{image_extension}')
            with open(image_path, 'wb') as img_file:
                img_file.write(image['image_data'])

    def save_tables(self, tables_data):
        """Save extracted tables as CSV files."""
        for i, table in enumerate(tables_data):
            # Safely get 'page_number', default to 'unknown_page' if missing
            page_number = table.get("page_number", "unknown_page")
            table_rows = table.get("table", [])
            table_path = os.path.join(self.output_directory, f'table_{i}_page_{page_number}.csv')
            
            with open(table_path, 'w', newline='') as csvfile:
                writer = csv.writer(csvfile)
                writer.writerows(table_rows)

    def save_links(self, links_data):
        """Save extracted hyperlinks to a text file."""
        with open(os.path.join(self.output_directory, 'extracted_links.txt'), 'w') as f:
            for link in links_data:
                # Use .get() to avoid KeyError if key is missing
                text = link.get('text', 'No text')  # Default to 'No text' if 'text' key is missing
                url = link.get('url', 'No URL')  # Default to 'No URL' if 'url' key is missing
                f.write(f"{text} -> {url}\n")

class MySQLStorage(Storage):
    def __init__(self, db_config):
        self.connection = mysql.connector.connect(**db_config)
        self.cursor = self.connection.cursor()
        self.create_tables()

    def create_tables(self):
        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS text_data (
                id INT AUTO_INCREMENT PRIMARY KEY,
                content TEXT NOT NULL,
                page_number INT
            )
        ''')

        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS images_data (
                id INT AUTO_INCREMENT PRIMARY KEY,
                image_data LONGBLOB NOT NULL,
                image_extension VARCHAR(10),
                page_number INT
            )
        ''')

        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS tables_data (
                id INT AUTO_INCREMENT PRIMARY KEY,
                table_data TEXT NOT NULL,
                page_number INT
            )
        ''')

        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS links_data (
                id INT AUTO_INCREMENT PRIMARY KEY,
                url TEXT NOT NULL,
                page_number INT
            )
        ''')
        self.connection.commit()

    def save_text(self, text_data):
        for item in text_data:
            self.cursor.execute('''
                INSERT INTO text_data (content, page_number) VALUES (%s, %s)
            ''', (item.get("text", ""), item.get("slide_number", None)))  # Use "slide_number" instead of "page_number"
        self.connection.commit()


    def save_images(self, images_data):
        for item in images_data:
            self.cursor.execute('''
                INSERT INTO images_data (image_data, image_extension, page_number) VALUES (%s, %s, %s)
            ''', (item["image_data"], item["image_extension"], item.get("page_number", None)))
        self.connection.commit()

    def save_tables(self, tables_data):
        for item in tables_data:
            self.cursor.execute('''
                INSERT INTO tables_data (table_data, page_number) VALUES (%s, %s)
            ''', (str(item["table"]), item.get("page_number", None)))
        self.connection.commit()

    def save_links(self, links_data):
        for item in links_data:
            self.cursor.execute('''
                INSERT INTO links_data (url, page_number) VALUES (%s, %s)
            ''', (item["url"], item.get("page_number", None)))
        self.connection.commit()

    def close(self):
        self.cursor.close()
        self.connection.close()                

class Processing:
    def process_file(loader_class, file_path, output_folder, db_config):
        loader = loader_class(file_path)
        extractor = DataExtractor(loader)

        # Extract data
        text_data = extractor.extract_text()
        link_data = extractor.extract_links()
        images_data = extractor.extract_images()
        tables_data = extractor.extract_tables()

        # Save to file storage
        file_storage = FileStorage(output_folder)
        file_storage.save_text(text_data)
        file_storage.save_links(link_data)
        file_storage.save_images(images_data)
        file_storage.save_tables(tables_data)

        # Save to MySQL storage
        mysql_storage = MySQLStorage(db_config)
        mysql_storage.save_text(text_data)
        mysql_storage.save_images(images_data)
        mysql_storage.save_tables(tables_data)
        mysql_storage.save_links(link_data)
        mysql_storage.close()


if __name__ == "__main__":
    # Take relative file path input from the user
    file_path = input("Enter the relative path to the file (with extension): ").strip()
    
    # Check if the relative path is valid
    if not os.path.isfile(file_path):
        print(f"The file at the path '{file_path}' does not exist. Please provide a valid relative path.")
    else:
        # Extract the file extension
        file_extension = file_path.split('.')[-1].lower() if '.' in file_path else ''

        # Map file extensions to their corresponding loader classes and output folders
        file_map = {
            'pdf': (PDFLoader, '/home/shtlp_0103/Assignment_3/Output/output1'),
            'docx': (DOCXLoader, '/home/shtlp_0103/Assignment_3/Output/output2'),
            'pptx': (PPTLoader, '/home/shtlp_0103/Assignment_3/Output/output3')
        }

        # Database configuration
        db_config = {
            'user': 'root',
            'password': 'santosh25',
            'host': 'localhost',
            'database': 'sql_storage',
        }

        # Check if the file extension is valid
        if file_extension in file_map:
            loader_class, output_folder = file_map[file_extension]
            Processing.process_file(loader_class, file_path, output_folder, db_config)
        else:
            print("Unsupported file type. Please enter a valid filename with a supported extension (pdf, docx, pptx).")