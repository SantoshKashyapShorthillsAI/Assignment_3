import random
from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()

# Example list of random hyperlinks, titles, and body texts
titles = [
    "Python Overview", "Data Science", "Machine Learning", "Artificial Intelligence",
    "Web Development", "Software Engineering", "Cybersecurity", "Cloud Computing",
    "Mobile App Development", "Game Development"
]
body_texts = [
    "This slide covers an introduction to the topic.", 
    "Here is some additional information about the subject.", 
    "You can find more details in the official documentation.", 
    "Check out the resources available online.", 
    "Stay tuned for the upcoming updates."
]
links = [
    "https://www.python.org", "https://www.djangoproject.com", "https://www.tensorflow.org", 
    "https://pytorch.org", "https://www.aws.amazon.com", "https://azure.microsoft.com", 
    "https://developer.android.com", "https://www.unrealengine.com"
]
link_texts = [
    "Visit Python's official website", "Check out Django", "Learn about TensorFlow",
    "Explore PyTorch", "Amazon Web Services", "Microsoft Azure", 
    "Android Development", "Unreal Engine"
]

# Add random slides
for i in range(50):
    slide_layout = presentation.slide_layouts[5]  # Use blank layout for simplicity
    slide = presentation.slides.add_slide(slide_layout)

    # Random title and body text
    title_text = random.choice(titles)
    body_text = random.choice(body_texts)
    
    # Add Title
    title = slide.shapes.title
    title.text = f"Slide {i+1}: {title_text}"

    # Add Body Text
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = body_text

    # Randomly add hyperlinks to some slides
    if random.choice([True, False]):  # 50% chance to add a hyperlink
        link = random.choice(links)
        link_text = random.choice(link_texts)
        p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = f"Click here to visit {link_text}"
        run.hyperlink.address = link

    # Randomly add a table to some slides
    if random.choice([True, False]):  # 50% chance to add a table
        rows = 2
        cols = 2
        left = Inches(1)
        top = Inches(3)
        width = Inches(5)
        height = Inches(1)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set table content
        table.cell(0, 0).text = 'Header 1'
        table.cell(0, 1).text = 'Header 2'
        table.cell(1, 0).text = 'Content 1'
        table.cell(1, 1).text = 'Content 2'

    # Randomly add images to some slides
    if random.choice([True, False]):  # 50% chance to add an image
        img_path = "/home/shtlp_0103/Assignment_3/Documents/apple.jpeg"  # Replace with valid path
        left = Inches(1)
        top = Inches(2)
        height = Inches(3)
        slide.shapes.add_picture(img_path, left, top, height=height)

# Save the presentation
presentation.save('large.pptx')
print("Presentation with 50+ slides created successfully!")
