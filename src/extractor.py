from abc import ABC, abstractmethod
import fitz
import docx
from pptx import Presentation
import os
import csv
import pdfplumber
import mysql.connector

class FileLoader(ABC):
    def __init__(self, file_path, expected_extension):
        self.file_path = file_path
        self.expected_extension = expected_extension

    def validate_extension(self):
        if not self.file_path.endswith(self.expected_extension):
            raise ValueError(f"Invalid file format. Expected {self.expected_extension}.")

    @abstractmethod
    def load(self):
        """Load the file content."""
        pass


class PDFLoader(FileLoader):
    def __init__(self, file_path):
        super().__init__(file_path, '.pdf')
        self.doc = None

    def load(self):
        self.validate_extension()
        self.doc = fitz.open(self.file_path)
        return self.doc


class DOCXLoader(FileLoader):
    def __init__(self, file_path):
        super().__init__(file_path, '.docx')
        self.doc = None

    def load(self):
        self.validate_extension()
        self.doc = docx.Document(self.file_path)
        return self.doc


class PPTLoader(FileLoader):
    def __init__(self, file_path):
        super().__init__(file_path, '.pptx')
        self.presentation = None

    def load(self):
        self.validate_extension()
        self.presentation = Presentation(self.file_path)
        return self.presentation

class DataExtractor:
    def __init__(self, file_loader: FileLoader):
        self.file_loader = file_loader
        self.file_loader.load()

    def extract_text(self):
        """Extracts text with metadata like page number and font details."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_text()
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_text()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_text()
        else:
            raise ValueError("Unsupported file type for text extraction.")

    def _extract_pdf_text(self):
        text_data = []
        for page_num in range(len(self.file_loader.doc)):
            page = self.file_loader.doc.load_page(page_num)
            text = page.get_text("text")
            text_data.append({
                "page_number": page_num + 1,
                "text": text
            })
        return text_data

    def _extract_docx_text(self):
        text_data = []
        for para in self.file_loader.doc.paragraphs:
            text_data.append({
                "text": para.text,
                "style": para.style.name
            })
        return text_data

    def _extract_ppt_text(self):
        text_data = []
        for slide_num, slide in enumerate(self.file_loader.presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            # Join the slide text list into a single string separated by newlines
            text_data.append({
                "slide_number": slide_num + 1,
                "text": "\n".join(slide_text)  # Join text list into a single string
            })
        return text_data

    def extract_links(self):
        """Extracts hyperlinks with metadata."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_links()
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_links()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_links()
        else:
            raise ValueError("Unsupported file type for link extraction.")

    def _extract_pdf_links(self):
        link_data = []
        for page_num in range(len(self.file_loader.doc)):
            page = self.file_loader.doc.load_page(page_num)
            links = page.get_links()
            for link in links:
                link_data.append({
                    "page_number": page_num + 1,
                    "url": link.get('uri')
                })
        return link_data

    def _extract_docx_links(self):
        link_data = []
        for para_num, para in enumerate(self.file_loader.doc.paragraphs):
            for run in para.runs:
                if run.font.color and run.text.startswith('http'):
                    link_data.append({
                        "paragraph_number": para_num + 1,  # Add paragraph number
                        "url": run.text,
                        "style": para.style.name
                    })
        return link_data


    def _extract_ppt_links(self):
        link_data = []
        for slide_num, slide in enumerate(self.file_loader.presentation.slides):
            for shape in slide.shapes:
                # Check if the shape contains text and has a hyperlink attribute
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink and run.hyperlink.address:
                                link_data.append({
                                    "slide_number": slide_num + 1,
                                    "url": run.hyperlink.address
                                })
                # In case the shape has a hyperlink directly (without being in the text frame)
                elif hasattr(shape, "hyperlink") and shape.hyperlink.address:
                    link_data.append({
                        "slide_number": slide_num + 1,
                        "url": shape.hyperlink.address
                    })
        return link_data


    def extract_images(self):
        """Extract images from the document."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_images()
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_images()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_images()
        else:
            raise ValueError("Unsupported file type for image extraction.")

    def _extract_pdf_images(self):
        image_data = []
        for page_num in range(len(self.file_loader.doc)):
            page = self.file_loader.doc.load_page(page_num)
            images = page.get_images(full=True)
            for img_index, img in enumerate(images):
                xref = img[0]
                image = self.file_loader.doc.extract_image(xref)
                img_bytes = image["image"]
                img_extension = image["ext"]
                image_data.append({
                    "page_number": page_num + 1,
                    "image_data": img_bytes,
                    "image_extension": img_extension
                })
        return image_data

    def _extract_docx_images(self):
        image_data = []
        for rel in self.file_loader.doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data.append({
                    "image_data": rel.target_part.blob,
                    "image_extension": rel.target_part.content_type.split('/')[-1]
                })
        return image_data

    def _extract_ppt_images(self):
        image_data = []
        for slide_num, slide in enumerate(self.file_loader.presentation.slides):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    img_bytes = shape.image.blob
                    image_extension = shape.image.content_type.split('/')[-1]
                    image_data.append({
                        "slide_number": slide_num + 1,
                        "image_data": img_bytes,
                        "image_extension": image_extension
                    })
        return image_data

    def extract_tables(self):
        """Extract tables from the document."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_tables_with_plumber()  # Using 
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_tables()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_tables()  
        else:
            raise ValueError("Unsupported file type for table extraction.")

    def _extract_pdf_tables_with_plumber(self):
        """Extract tables from PDF using pdfplumber."""
        table_data = []
        with pdfplumber.open(self.file_loader.file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                tables = page.extract_tables()  # Extract tables from each page
                for table in tables:
                    table_data.append({
                        "page_number": page_num + 1,
                        "table": table
                    })
        return table_data
    
    def _extract_docx_tables(self):
        table_data = []
        for table in self.file_loader.doc.tables:
            table_rows = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                table_rows.append(row_data)
            table_data.append({
                "table": table_rows
            })
        return table_data

    def _extract_ppt_tables(self):
        table_data = []
        for slide_num, slide in enumerate(self.file_loader.presentation.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_rows.append(row_data)
                    table_data.append({
                        "slide_number": slide_num + 1,
                        "table": table_rows
                    })
        return table_data

class Storage(ABC):
    """Abstract class for storing extracted data."""

    @abstractmethod
    def save_text(self, text_data):
        """Save extracted text."""
        pass

    @abstractmethod
    def save_images(self, images_data):
        """Save extracted images."""
        pass

    @abstractmethod
    def save_tables(self, tables_data):
        """Save extracted tables."""
        pass

    @abstractmethod
    def save_links(self, links_data):
        """Save extracted links."""
        pass


class FileStorage(Storage):
    """Concrete class for storing extracted data to files."""

    def __init__(self, output_directory):
        """Initialize FileStorage with an output directory."""
        self.output_directory = output_directory
        if not os.path.exists(self.output_directory):
            os.makedirs(self.output_directory)

    def save_text(self, text_data):
        """Save extracted text to a text file."""
        with open(os.path.join(self.output_directory, 'extracted_text.txt'), 'w') as f:
            for entry in text_data:
                f.write(f"{entry}\n")

    def save_links(self, links_data):
        """Save extracted hyperlinks with page/slide/paragraph number to a text file."""
        with open(os.path.join(self.output_directory, 'extracted_links.txt'), 'w') as f:
            for link in links_data:
                # Determine whether it's a PDF page, PowerPoint slide, or Word paragraph
                location = ""
                if 'page_number' in link:
                    location = f"Page {link['page_number']}"  # For PDF
                elif 'slide_number' in link:
                    location = f"Slide {link['slide_number']}"  # For PowerPoint
                elif 'paragraph_number' in link:
                    location = f"Paragraph {link['paragraph_number']}"  # For Word

                # Use .get() to avoid KeyError if 'url' is missing
                url = link.get('url', 'No URL')  # Default to 'No URL' if 'url' key is missing

                # Write the location (page/slide/paragraph) and URL to the file
                f.write(f"{location} -> {url}\n")


    def save_images(self, images_data):
        """Save extracted images to the output directory."""
        for i, image in enumerate(images_data):
            image_extension = image.get("image_extension", "png")  # Default to PNG if not provided
            image_path = os.path.join(self.output_directory, f'image_{i}.{image_extension}')
            with open(image_path, 'wb') as img_file:
                img_file.write(image['image_data'])



    def save_tables(self, tables_data):
        """Save extracted tables as CSV files along with metadata."""
        for i, table in enumerate(tables_data):
            # Safely get 'page_number', 'slide_number', or use 'unknown' if missing
            page_number = table.get("page_number", table.get("slide_number", "unknown_location"))
            
            # Safely get table data
            table_rows = table.get("table", [])

            # Define paths for the CSV file and metadata file
            table_path = os.path.join(self.output_directory, f'table_{i}_location_{page_number}.csv')
            metadata_path = os.path.join(self.output_directory, f'table_{i}_location_{page_number}_metadata.txt')

            # Save the table data as a CSV file
            with open(table_path, 'w', newline='') as csvfile:
                writer = csv.writer(csvfile)
                writer.writerows(table_rows)

            # Save metadata as a separate .txt file
            with open(metadata_path, 'w') as metafile:
                metafile.write(f"Table {i + 1} Metadata\n")
                if 'page_number' in table:
                    metafile.write(f"Extracted from PDF - Page {table['page_number']}\n")
                elif 'slide_number' in table:
                    metafile.write(f"Extracted from PowerPoint - Slide {table['slide_number']}\n")
                else:
                    metafile.write("Extracted from Word document\n")

                # You can add more metadata if needed (e.g., source, timestamp, etc.)
                metafile.write(f"Number of rows: {len(table_rows)}\n")
                if table_rows:
                    metafile.write(f"Number of columns: {len(table_rows[0])}\n")
                else:
                    metafile.write("Number of columns: 0\n")

class MySQLStorage(Storage):
    def __init__(self, db_config):
        self.connection = mysql.connector.connect(**db_config)
        self.cursor = self.connection.cursor()
        self.create_tables()

    def create_tables(self):
        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS text_data (
                id INT AUTO_INCREMENT PRIMARY KEY,
                content TEXT NOT NULL,
                page_number INT
            )
        ''')

        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS images_data (
                id INT AUTO_INCREMENT PRIMARY KEY,
                image_data LONGBLOB NOT NULL,
                image_extension VARCHAR(10),
                page_number INT
            )
        ''')

        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS tables_data (
                id INT AUTO_INCREMENT PRIMARY KEY,
                table_data TEXT NOT NULL,
                page_number INT
            )
        ''')

        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS links_data (
                id INT AUTO_INCREMENT PRIMARY KEY,
                url TEXT NOT NULL,
                page_number INT
            )
        ''')
        self.connection.commit()

    def save_text(self, text_data):
        for item in text_data:
            self.cursor.execute('''
                INSERT INTO text_data (content, page_number) VALUES (%s, %s)
            ''', (item.get("text", ""), item.get("slide_number", None))) 
        self.connection.commit()


    def save_images(self, images_data):
        for item in images_data:
            self.cursor.execute('''
                INSERT INTO images_data (image_data, image_extension, page_number) VALUES (%s, %s, %s)
            ''', (item["image_data"], item["image_extension"], item.get("page_number", None)))
        self.connection.commit()

    def save_tables(self, tables_data):
        for item in tables_data:
            self.cursor.execute('''
                INSERT INTO tables_data (table_data, page_number) VALUES (%s, %s)
            ''', (str(item["table"]), item.get("page_number", None)))
        self.connection.commit()

    def save_links(self, links_data):
        for item in links_data:
            url = item.get("url")
            if url:  # Ensure url is not None or empty
                self.cursor.execute('''
                    INSERT INTO links_data (url, page_number) VALUES (%s, %s)
                ''', (url, item.get("page_number", None)))
        self.connection.commit()

    def close(self):
        self.cursor.close()
        self.connection.close()                

class Processing:
    def process_file(loader_class, file_path, output_folder, db_config):
        loader = loader_class(file_path)
        extractor = DataExtractor(loader)

        # Extract data
        text_data = extractor.extract_text()
        link_data = extractor.extract_links()
        images_data = extractor.extract_images()
        tables_data = extractor.extract_tables()

        # Save to file storage
        file_storage = FileStorage(output_folder)
        file_storage.save_text(text_data)
        file_storage.save_links(link_data)
        file_storage.save_images(images_data)
        file_storage.save_tables(tables_data)

        # Save to MySQL storage
        mysql_storage = MySQLStorage(db_config)
        mysql_storage.save_text(text_data)
        mysql_storage.save_images(images_data)
        mysql_storage.save_tables(tables_data)
        mysql_storage.save_links(link_data)
        mysql_storage.close()
        

if __name__ == "__main__":
    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    # Directories
    base_dir = os.path.join(project_root, "Documents")
    output_dir = os.path.join(project_root, "Output")
    
    # Print base directory for debugging
    print(f"Base Directory: {base_dir}")

    # Get filename from user
    file_name = input("Enter the filename (with extension): ").strip()
    file_path = os.path.join(base_dir, file_name)

    # Check if file exists
    if not os.path.isfile(file_path):
        print(f"The file at the path '{file_path}' does not exist. Please provide a valid relative path.")
    else:
        # Map file extensions to loader classes and output folders
        file_map = {
            'pdf': (PDFLoader, os.path.join(output_dir, "PDF")),
            'docx': (DOCXLoader, os.path.join(output_dir, "DOCX")),
            'pptx': (PPTLoader, os.path.join(output_dir, "PPTX")),
        }

        # Extract file extension and process if valid
        file_extension = file_path.split('.')[-1].lower()
        loader_class_output = file_map.get(file_extension)
        
        if loader_class_output:
            loader_class, output_folder = loader_class_output
            db_config = {
                'user': os.getenv('DB_USER'),
                'password': os.getenv('DB_PASSWORD'),
                'host': os.getenv('DB_HOST'),
                'database': os.getenv('DB_DATABASE'),
            }
            Processing.process_file(loader_class, file_path, output_folder, db_config)
        else:
            print("Unsupported file type. Please enter a valid filename with a supported extension (pdf, docx, pptx).")
