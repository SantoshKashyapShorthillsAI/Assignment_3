# data_extractors.py

from abc import ABC, abstractmethod
from file_loaders import PDFLoader, DOCXLoader, PPTLoader, FileLoader
import os
import fitz
import docx
from pptx import Presentation
import pdfplumber
import logging

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

class DataExtractor:
    """
    Class for extracting data from different file types (PDF, DOCX, PPTX).

    Attributes:
        file_loader (FileLoader): An instance of a file loader for loading specific file types.

    Methods:
        extract_text(): Extracts text with metadata (page number, style).
        extract_links(): Extracts hyperlinks with metadata.
        extract_images(): Extracts images from the document.
        extract_tables(): Extracts tables from the document.
    """

    def __init__(self, file_loader: FileLoader):
        """
        Initializes the DataExtractor.

        Args:
            file_loader (FileLoader): An instance of the FileLoader (PDFLoader, DOCXLoader, or PPTLoader).
        """
        self.file_loader = file_loader
        try:
            self.file_loader.load()
        except Exception as e:
            logging.error(f"Failed to load the file: {str(e)}")
            raise RuntimeError(f"Failed to load the file: {str(e)}")

    def extract_text(self):
        """Extracts text with metadata like page number and font details."""
        return self._extract_generic(self._extract_text_for_loader)

    def extract_links(self):
        """Extracts hyperlinks with metadata."""
        return self._extract_generic(self._extract_links_for_loader)

    def extract_images(self):
        """Extract images from the document."""
        return self._extract_generic(self._extract_images_for_loader)

    def extract_tables(self):
        """Extract tables from the document."""
        return self._extract_generic(self._extract_tables_for_loader)

    def _extract_generic(self, extractor_method):
        """Generic extraction method to handle the extraction based on loader type."""
        try:
            return extractor_method()
        except Exception as e:
            logging.error(f"Error during extraction: {str(e)}")
            raise RuntimeError(f"Error during extraction: {str(e)}")

    def _extract_text_for_loader(self):
        """Determines which text extraction method to call based on the loader type."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_text()
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_text()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_text()
        else:
            raise ValueError("Unsupported file type for text extraction.")

    def _extract_pdf_text(self):
        """Extracts text from a PDF file."""
        text_data = []
        try:
            for page_num in range(len(self.file_loader.doc)):
                page = self.file_loader.doc.load_page(page_num)
                text = page.get_text("text")
                text_data.append({
                    "page_number": page_num + 1,
                    "text": text
                })
        except Exception as e:
            logging.error(f"Error extracting text from PDF: {str(e)}")
            raise RuntimeError(f"Error extracting text from PDF: {str(e)}")
        return text_data

    def _extract_docx_text(self):
        """Extracts text from a DOCX file."""
        text_data = []
        try:
            for para in self.file_loader.doc.paragraphs:
                text_data.append({
                    "text": para.text,
                    "style": para.style.name
                })
        except Exception as e:
            logging.error(f"Error extracting text from DOCX: {str(e)}")
            raise RuntimeError(f"Error extracting text from DOCX: {str(e)}")
        return text_data

    def _extract_ppt_text(self):
        """Extracts text from a PPTX file."""
        text_data = []
        try:
            for slide_num, slide in enumerate(self.file_loader.presentation.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                # Join the slide text list into a single string separated by newlines
                text_data.append({
                    "slide_number": slide_num + 1,
                    "text": "\n".join(slide_text)  # Join text list into a single string
                })
        except Exception as e:
            logging.error(f"Error extracting text from PPTX: {str(e)}")
            raise RuntimeError(f"Error extracting text from PPTX: {str(e)}")
        return text_data

    def _extract_links_for_loader(self):
        """Determines which link extraction method to call based on the loader type."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_links()
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_links()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_links()
        else:
            raise ValueError("Unsupported file type for link extraction.")

    def _extract_pdf_links(self):
        """Extracts hyperlinks from a PDF file."""
        link_data = []
        try:
            for page_num in range(len(self.file_loader.doc)):
                page = self.file_loader.doc.load_page(page_num)
                links = page.get_links()
                for link in links:
                    link_data.append({
                        "page_number": page_num + 1,
                        "url": link.get('uri')
                    })
        except Exception as e:
            logging.error(f"Error extracting links from PDF: {str(e)}")
            raise RuntimeError(f"Error extracting links from PDF: {str(e)}")
        return link_data

    def _extract_docx_links(self):
        """Extracts hyperlinks from a DOCX file."""
        link_data = []
        try:
            for para_num, para in enumerate(self.file_loader.doc.paragraphs):
                for run in para.runs:
                    if run.font.color and run.text.startswith('http'):
                        link_data.append({
                            "paragraph_number": para_num + 1,  # Add paragraph number
                            "url": run.text,
                            "style": para.style.name
                        })
        except Exception as e:
            logging.error(f"Error extracting links from DOCX: {str(e)}")
            raise RuntimeError(f"Error extracting links from DOCX: {str(e)}")
        return link_data

    def _extract_ppt_links(self):
        """Extracts hyperlinks from a PPTX file."""
        link_data = []
        try:
            for slide_num, slide in enumerate(self.file_loader.presentation.slides):
                for shape in slide.shapes:
                    # Check if the shape contains text and has a hyperlink attribute
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.hyperlink and run.hyperlink.address:
                                    link_data.append({
                                        "slide_number": slide_num + 1,
                                        "url": run.hyperlink.address
                                    })
                    elif hasattr(shape, "hyperlink") and shape.hyperlink.address:
                        link_data.append({
                            "slide_number": slide_num + 1,
                            "url": shape.hyperlink.address
                        })
        except Exception as e:
            logging.error(f"Error extracting links from PPTX: {str(e)}")
            raise RuntimeError(f"Error extracting links from PPTX: {str(e)}")
        return link_data

    def _extract_images_for_loader(self):
        """Determines which image extraction method to call based on the loader type."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_images()
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_images()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_images()
        else:
            raise ValueError("Unsupported file type for image extraction.")

    def _extract_pdf_images(self):
        """Extracts images from a PDF file."""
        image_data = []
        try:
            for page_num in range(len(self.file_loader.doc)):
                page = self.file_loader.doc.load_page(page_num)
                images = page.get_images(full=True)
                for img_index, img in enumerate(images):
                    xref = img[0]
                    image = self.file_loader.doc.extract_image(xref)
                    img_bytes = image["image"]
                    img_extension = image["ext"]
                    image_data.append({
                        "page_number": page_num + 1,
                        "image_data": img_bytes,
                        "image_extension": img_extension
                    })
        except Exception as e:
            logging.error(f"Error extracting images from PDF: {str(e)}")
            raise RuntimeError(f"Error extracting images from PDF: {str(e)}")
        return image_data

    def _extract_docx_images(self):
        """Extracts images from a DOCX file."""
        image_data = []
        try:
            for rel in self.file_loader.doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_data.append({
                        "image_data": rel.target_part.blob,
                        "image_extension": rel.target_part.content_type.split('/')[-1]
                    })
        except Exception as e:
            logging.error(f"Error extracting images from DOCX: {str(e)}")
            raise RuntimeError(f"Error extracting images from DOCX: {str(e)}")
        return image_data

    def _extract_ppt_images(self):
        """Extracts images from a PPTX file."""
        image_data = []
        try:
            for slide_num, slide in enumerate(self.file_loader.presentation.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Shape type for pictures
                        img_stream = shape.image.blob
                        img_extension = shape.image.ext
                        image_data.append({
                            "slide_number": slide_num + 1,
                            "image_data": img_stream,
                            "image_extension": img_extension
                        })
        except Exception as e:
            logging.error(f"Error extracting images from PPTX: {str(e)}")
            raise RuntimeError(f"Error extracting images from PPTX: {str(e)}")
        return image_data

    def _extract_tables_for_loader(self):
        """Determines which table extraction method to call based on the loader type."""
        if isinstance(self.file_loader, PDFLoader):
            return self._extract_pdf_tables()
        elif isinstance(self.file_loader, DOCXLoader):
            return self._extract_docx_tables()
        elif isinstance(self.file_loader, PPTLoader):
            return self._extract_ppt_tables()
        else:
            raise ValueError("Unsupported file type for table extraction.")

    def _extract_pdf_tables(self):
        """Extracts tables from a PDF file using pdfplumber."""
        table_data = []
        try:
            with pdfplumber.open(self.file_loader.file_path) as pdf:
                for page_num in range(len(pdf.pages)):
                    page = pdf.pages[page_num]
                    tables = page.extract_tables()
                    for table in tables:
                        table_data.append({
                            "page_number": page_num + 1,
                            "table": table
                        })
        except Exception as e:
            logging.error(f"Error extracting tables from PDF: {str(e)}")
            raise RuntimeError(f"Error extracting tables from PDF: {str(e)}")
        return table_data

    def _extract_docx_tables(self):
        """Extracts tables from a DOCX file."""
        table_data = []
        try:
            for table_num, table in enumerate(self.file_loader.doc.tables):
                rows = []
                for row in table.rows:
                    cols = [cell.text for cell in row.cells]
                    rows.append(cols)
                table_data.append({
                    "table_number": table_num + 1,
                    "table": rows
                })
        except Exception as e:
            logging.error(f"Error extracting tables from DOCX: {str(e)}")
            raise RuntimeError(f"Error extracting tables from DOCX: {str(e)}")
        return table_data

    def _extract_ppt_tables(self):
        """Extracts tables from a PPTX file."""
        table_data = []
        try:
            for slide_num, slide in enumerate(self.file_loader.presentation.slides):
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        rows = []
                        for row in table.rows:
                            cols = [cell.text for cell in row.cells]
                            rows.append(cols)
                        table_data.append({
                            "slide_number": slide_num + 1,
                            "table": rows
                        })
        except Exception as e:
            logging.error(f"Error extracting tables from PPTX: {str(e)}")
            raise RuntimeError(f"Error extracting tables from PPTX: {str(e)}")
        return table_data
