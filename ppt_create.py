from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()

# Add a slide for text and links
slide_layout = presentation.slide_layouts[5]  # 5 is a blank slide
slide = presentation.slides.add_slide(slide_layout)

# Add Title Text
title = slide.shapes.title
title.text = "PowerPoint with Links, Text, Table, and Images"

# Add body text box
left = Inches(1)
top = Inches(2)
width = Inches(6)
height = Inches(1)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = "This is a Python-generated PowerPoint with multiple elements."

# Add a hyperlink to the text
p = text_frame.add_paragraph()
run = p.add_run()
run.text = "Click here to visit Python's official website"
run.hyperlink.address = "https://www.python.org"

# Add a slide for a table
slide_layout = presentation.slide_layouts[5]
slide2 = presentation.slides.add_slide(slide_layout)

# Add a title for the table slide
title = slide2.shapes.title
title.text = "Table Example"

# Create a table
rows = 2
cols = 2
left = Inches(2)
top = Inches(2)
width = Inches(4)
height = Inches(1)
table = slide2.shapes.add_table(rows, cols, left, top, width, height).table

# Set column headers
table.cell(0, 0).text = 'Name'
table.cell(0, 1).text = 'Age'

# Set row content
table.cell(1, 0).text = 'John'
table.cell(1, 1).text = '30'

# Add a slide for an image
slide3 = presentation.slides.add_slide(slide_layout)

# Add a title for the image slide
title = slide3.shapes.title
title.text = "Image Example"

# Add an image (Make sure you have an image in the same directory as the script)
img_path = "apple.jpeg"  # Update with your image path
left = Inches(1)
top = Inches(2)
height = Inches(3)
slide3.shapes.add_picture(img_path, left, top, height=height)

# Save the presentation
presentation.save('ppt_example.pptx')
print("Presentation created successfully!")
